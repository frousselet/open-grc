"""Management review report generators (PPTX presentation and DOCX minutes).

Sections follow ISO 27001:2022 clause 9.3 management review inputs:
1. Status of actions from previous management reviews
2. Changes to external and internal issues relevant to the ISMS
3. Changes to needs and expectations of interested parties
4. Feedback on information security performance
   a) Non-conformities and corrective actions
   b) Monitoring and measurement results
   c) Audit results
   d) Achievement of information security objectives
5. Feedback from interested parties
6. Risk assessment results and risk treatment plan status
7. Opportunities for improvement
"""

import io
import re
from collections import Counter

from django.utils import timezone
from django.utils.html import strip_tags

from accounts.models import CompanySettings

# ── Shared colour palette (neutral, consistent with audit_report_pdf) ──
_CLR_DARK = (0x22, 0x22, 0x22)           # body text
_CLR_HEADING = (0x33, 0x33, 0x33)        # headings
_CLR_HEADER_BG = (0xDC, 0xE6, 0xF0)     # table header background (light steel)
_CLR_HEADER_FG = (0x22, 0x22, 0x22)      # table header text
_CLR_ALT_ROW = (0xF5, 0xF7, 0xFA)       # alternating row
_CLR_BORDER = (0xBB, 0xBB, 0xBB)        # table border
_CLR_ACCENT = (0x33, 0x6B, 0x9B)        # subtle accent (section dividers)
_CLR_MUTED = (0x77, 0x77, 0x77)         # secondary text


def _strip(text):
    """Strip HTML and clean whitespace."""
    if not text:
        return ""
    cleaned = re.sub(r"<p>\s*(?:&nbsp;|<br\s*/?>)?\s*</p>", "", text)
    return strip_tags(cleaned).strip()


# ═══════════════════════════════════════════════════════════════════
# Data gathering (shared between PPTX and DOCX)
# ═══════════════════════════════════════════════════════════════════


def gather_management_review_data(user, scope_ids=None,
                                   period_start=None, period_end=None):
    """Gather all data needed for the management review report.

    Args:
        period_start: Optional start date of the review period.
        period_end: Optional end date of the review period (defaults to today).

    Returns a dict with all sections pre-computed.
    """
    from datetime import datetime

    from compliance.constants import (
        ActionPlanStatus,
        AssessmentStatus,
        FindingType,
    )
    from compliance.models import (
        ComplianceActionPlan,
        ComplianceAssessment,
        Finding,
        Framework,
    )
    from context.constants import ObjectiveStatus
    from context.models import (
        Indicator,
        Issue,
        Objective,
        Stakeholder,
        StakeholderExpectation,
    )
    from risks.models import Risk, RiskAssessment, RiskTreatmentPlan

    now = timezone.now()
    company = CompanySettings.get()

    # Convert period dates to timezone-aware datetimes for filtering
    p_start_dt = (
        timezone.make_aware(datetime.combine(period_start, datetime.min.time()))
        if period_start else None
    )
    p_end_dt = (
        timezone.make_aware(datetime.combine(period_end, datetime.max.time()))
        if period_end else None
    )

    # Helper: apply period filter on a datetime field
    def _period_filter(qs, dt_field="updated_at"):
        if p_start_dt:
            qs = qs.filter(**{f"{dt_field}__gte": p_start_dt})
        if p_end_dt:
            qs = qs.filter(**{f"{dt_field}__lte": p_end_dt})
        return qs

    def _period_filter_date(qs, date_field):
        if period_start:
            qs = qs.filter(**{f"{date_field}__gte": period_start})
        if period_end:
            qs = qs.filter(**{f"{date_field}__lte": period_end})
        return qs

    # ── 1. Status of actions from previous management reviews ──
    action_plans = ComplianceActionPlan.objects.select_related(
        "owner",
    ).prefetch_related("requirements", "findings", "risks")
    if scope_ids:
        action_plans = action_plans.filter(scopes__id__in=scope_ids)
    action_plans = _period_filter(action_plans)

    action_plan_stats = {
        "total": action_plans.count(),
        "by_status": {},
        "overdue": action_plans.filter(
            target_date__lt=now.date(),
        ).exclude(
            status__in=[
                ActionPlanStatus.CLOSED,
                ActionPlanStatus.CANCELLED,
                ActionPlanStatus.VALIDATED,
            ],
        ).count(),
    }
    for status_value, status_label in ActionPlanStatus.choices:
        count = action_plans.filter(status=status_value).count()
        if count:
            action_plan_stats["by_status"][str(status_label)] = count

    active_action_plans = list(
        action_plans.exclude(
            status__in=[ActionPlanStatus.CLOSED, ActionPlanStatus.CANCELLED],
        ).order_by("target_date")[:20]
    )
    action_plan_rows = []
    for ap in active_action_plans:
        action_plan_rows.append({
            "reference": ap.reference,
            "name": ap.name,
            "status": ap.get_status_display(),
            "owner": str(ap.owner) if ap.owner else "-",
            "target_date": ap.target_date.strftime("%d/%m/%Y") if ap.target_date else "-",
            "progress": f"{ap.progress_percentage}%" if ap.progress_percentage is not None else "-",
            "is_overdue": (
                ap.target_date
                and ap.target_date < now.date()
                and ap.status not in [
                    ActionPlanStatus.CLOSED,
                    ActionPlanStatus.CANCELLED,
                    ActionPlanStatus.VALIDATED,
                ]
            ),
        })

    # ── 2. Changes to external and internal issues ──
    issues = Issue.objects.all()
    if scope_ids:
        issues = issues.filter(scopes__id__in=scope_ids)
    issues = _period_filter(issues)

    internal_issues = list(issues.filter(type="internal").order_by("-updated_at")[:15])
    external_issues = list(issues.filter(type="external").order_by("-updated_at")[:15])

    def _issue_row(issue):
        return {
            "reference": issue.reference,
            "name": issue.name,
            "type": issue.get_type_display(),
            "category": issue.get_category_display(),
            "impact_level": issue.get_impact_level_display(),
            "trend": issue.get_trend_display() if issue.trend else "-",
            "status": issue.get_status_display(),
        }

    # ── 3. Changes to needs and expectations of interested parties ──
    stakeholders = Stakeholder.objects.prefetch_related("expectations")
    if scope_ids:
        stakeholders = stakeholders.filter(scopes__id__in=scope_ids)
    stakeholders = _period_filter(stakeholders)
    stakeholders = list(stakeholders.order_by("-updated_at")[:20])

    stakeholder_rows = []
    for sh in stakeholders:
        expectations = list(sh.expectations.filter(is_applicable=True))
        stakeholder_rows.append({
            "reference": sh.reference,
            "name": sh.name,
            "type": sh.get_type_display(),
            "category": sh.get_category_display(),
            "influence_level": sh.get_influence_level_display(),
            "interest_level": sh.get_interest_level_display(),
            "expectations_count": len(expectations),
            "top_expectations": [
                _strip(e.description) for e in expectations[:3]
            ],
        })

    # ── 4a. Non-conformities and corrective actions ──
    findings = Finding.objects.select_related("assessment", "assessor")
    findings = _period_filter(findings, "created_at")
    nc_findings = list(
        findings.filter(
            finding_type__in=[FindingType.MAJOR_NON_CONFORMITY, FindingType.MINOR_NON_CONFORMITY],
        ).order_by("-created_at")[:20]
    )
    finding_type_counts = Counter(
        findings.values_list("finding_type", flat=True)
    )
    nc_rows = []
    for f in nc_findings:
        nc_rows.append({
            "reference": f.reference,
            "type": f.get_finding_type_display(),
            "description": _strip(f.description)[:200],
            "assessment": f.assessment.name if f.assessment else "-",
            "recommendation": _strip(f.recommendation)[:200],
        })

    # ── 4b. Monitoring and measurement results ──
    indicators = Indicator.objects.all()
    if scope_ids:
        indicators = indicators.filter(scopes__id__in=scope_ids)
    indicators = list(indicators.filter(status="active").order_by("name")[:20])

    indicator_rows = []
    for ind in indicators:
        indicator_rows.append({
            "reference": ind.reference,
            "name": ind.name,
            "current_value": ind.current_value or "-",
            "expected_level": ind.expected_level or "-",
            "unit": ind.unit,
            "is_critical": ind.is_critical,
            "review_frequency": ind.get_review_frequency_display(),
        })

    # ── 4c. Audit results ──
    assessments = ComplianceAssessment.objects.prefetch_related("frameworks")
    if scope_ids:
        assessments = assessments.filter(scopes__id__in=scope_ids)
    assessments = _period_filter_date(assessments, "assessment_end_date")
    recent_assessments = list(
        assessments.filter(
            status__in=[AssessmentStatus.COMPLETED, AssessmentStatus.CLOSED],
        ).order_by("-assessment_end_date")[:10]
    )

    assessment_rows = []
    for a in recent_assessments:
        fw_names = ", ".join(fw.short_name or fw.name for fw in a.frameworks.all())
        assessment_rows.append({
            "reference": a.reference,
            "name": a.name,
            "frameworks": fw_names,
            "status": a.get_status_display(),
            "compliance_level": f"{a.overall_compliance_level:.0f}%" if a.overall_compliance_level is not None else "-",
            "coverage": f"{a.coverage_pct:.0f}%" if a.coverage_pct is not None else "-",
            "end_date": a.assessment_end_date.strftime("%d/%m/%Y") if a.assessment_end_date else "-",
        })

    # ── 4d. Achievement of information security objectives ──
    objectives = Objective.objects.select_related("owner")
    if scope_ids:
        objectives = objectives.filter(scopes__id__in=scope_ids)
    objectives = _period_filter(objectives)
    objectives = list(objectives.filter(
        status__in=[ObjectiveStatus.ACTIVE, ObjectiveStatus.ACHIEVED, ObjectiveStatus.NOT_ACHIEVED],
    ).order_by("status", "name")[:20])

    objective_rows = []
    for obj in objectives:
        objective_rows.append({
            "reference": obj.reference,
            "name": obj.name,
            "category": obj.get_category_display(),
            "status": obj.get_status_display(),
            "progress": f"{obj.progress_percentage}%" if obj.progress_percentage is not None else "-",
            "target_value": obj.target_value or "-",
            "current_value": obj.current_value or "-",
            "owner": str(obj.owner) if obj.owner else "-",
            "target_date": obj.target_date.strftime("%d/%m/%Y") if obj.target_date else "-",
        })

    objective_stats = {
        "total": len(objectives),
        "achieved": sum(1 for o in objectives if o.status == ObjectiveStatus.ACHIEVED),
        "not_achieved": sum(1 for o in objectives if o.status == ObjectiveStatus.NOT_ACHIEVED),
        "in_progress": sum(1 for o in objectives if o.status == ObjectiveStatus.ACTIVE),
    }

    # ── 5. Feedback from interested parties ──
    expectations = StakeholderExpectation.objects.select_related(
        "stakeholder",
    ).filter(is_applicable=True)
    expectations = _period_filter(expectations)
    expectations = expectations.order_by("-updated_at")[:20]
    expectation_rows = []
    for e in expectations:
        expectation_rows.append({
            "stakeholder": e.stakeholder.name,
            "description": _strip(e.description)[:200],
            "type": e.get_type_display(),
            "priority": e.get_priority_display(),
        })

    # ── 6. Risk assessment results and treatment plan status ──
    risk_assessments = RiskAssessment.objects.all()
    if scope_ids:
        risk_assessments = risk_assessments.filter(scopes__id__in=scope_ids)
    risk_assessments = _period_filter(risk_assessments, "created_at")
    risk_assessments = list(risk_assessments.order_by("-created_at")[:10])

    risk_assessment_rows = []
    for ra in risk_assessments:
        risk_count = ra.risks.count()
        risk_assessment_rows.append({
            "reference": ra.reference,
            "name": ra.name,
            "methodology": ra.get_methodology_display(),
            "status": ra.get_status_display(),
            "risk_count": risk_count,
            "date": ra.assessment_date.strftime("%d/%m/%Y") if ra.assessment_date else "-",
        })

    risks = Risk.objects.select_related("assessment")
    if scope_ids:
        risks = risks.filter(assessment__scopes__id__in=scope_ids)
    risks = _period_filter(risks)

    risk_stats = {
        "total": risks.count(),
        "by_treatment": {},
    }
    for decision_value, decision_label in [
        ("mitigate", "Mitigate"), ("accept", "Accept"),
        ("transfer", "Transfer"), ("avoid", "Avoid"),
        ("not_decided", "Not decided"),
    ]:
        count = risks.filter(treatment_decision=decision_value).count()
        if count:
            risk_stats["by_treatment"][decision_label] = count

    # Critical/high risks
    critical_risks = list(
        risks.filter(
            current_risk_level__gte=4,
        ).order_by("-current_risk_level")[:10]
    )
    critical_risk_rows = []
    for r in critical_risks:
        critical_risk_rows.append({
            "reference": r.reference,
            "name": r.name,
            "risk_level": r.current_risk_level or "-",
            "treatment": r.get_treatment_decision_display() if r.treatment_decision else "-",
            "status": r.get_status_display(),
        })

    treatment_plans = RiskTreatmentPlan.objects.select_related("risk", "owner")
    treatment_plans = _period_filter(treatment_plans)
    tp_stats = {
        "total": treatment_plans.count(),
        "by_status": {},
        "overdue": treatment_plans.filter(
            target_date__lt=now.date(),
        ).exclude(
            status__in=["completed", "cancelled"],
        ).count(),
    }
    for status_value, status_label in [
        ("planned", "Planned"), ("in_progress", "In progress"),
        ("completed", "Completed"), ("cancelled", "Cancelled"),
        ("overdue", "Overdue"),
    ]:
        count = treatment_plans.filter(status=status_value).count()
        if count:
            tp_stats["by_status"][str(status_label)] = count

    active_treatment_plans = list(
        treatment_plans.exclude(
            status__in=["completed", "cancelled"],
        ).order_by("target_date")[:15]
    )
    treatment_plan_rows = []
    for tp in active_treatment_plans:
        treatment_plan_rows.append({
            "reference": tp.reference,
            "name": tp.name,
            "risk": tp.risk.name if tp.risk else "-",
            "status": tp.get_status_display(),
            "owner": str(tp.owner) if tp.owner else "-",
            "target_date": tp.target_date.strftime("%d/%m/%Y") if tp.target_date else "-",
            "progress": f"{tp.progress_percentage}%",
        })

    # ── 7. Opportunities for improvement ──
    improvement_findings = list(
        findings.filter(
            finding_type=FindingType.IMPROVEMENT_OPPORTUNITY,
        ).order_by("-created_at")[:15]
    )
    improvement_rows = []
    for f in improvement_findings:
        improvement_rows.append({
            "reference": f.reference,
            "description": _strip(f.description)[:200],
            "recommendation": _strip(f.recommendation)[:200],
            "assessment": f.assessment.name if f.assessment else "-",
        })

    # ── Framework compliance summary ──
    frameworks = Framework.objects.all()
    framework_rows = []
    for fw in frameworks:
        level = fw.compliance_level
        framework_rows.append({
            "name": fw.short_name or fw.name,
            "compliance_level": f"{level:.0f}%" if level is not None else "-",
            "requirement_count": fw.requirements.count(),
            "status": fw.get_status_display(),
        })

    return {
        "company": company,
        "generated_at": now,
        "generated_by": user,
        "period_start": period_start,
        "period_end": period_end,
        # Section 1
        "action_plan_stats": action_plan_stats,
        "action_plan_rows": action_plan_rows,
        # Section 2
        "internal_issues": [_issue_row(i) for i in internal_issues],
        "external_issues": [_issue_row(i) for i in external_issues],
        # Section 3
        "stakeholder_rows": stakeholder_rows,
        # Section 4a
        "nc_rows": nc_rows,
        "finding_type_counts": {
            FindingType(k).label: v for k, v in finding_type_counts.items()
        } if finding_type_counts else {},
        # Section 4b
        "indicator_rows": indicator_rows,
        # Section 4c
        "assessment_rows": assessment_rows,
        # Section 4d
        "objective_rows": objective_rows,
        "objective_stats": objective_stats,
        # Section 5
        "expectation_rows": expectation_rows,
        # Section 6
        "risk_assessment_rows": risk_assessment_rows,
        "risk_stats": risk_stats,
        "critical_risk_rows": critical_risk_rows,
        "tp_stats": tp_stats,
        "treatment_plan_rows": treatment_plan_rows,
        # Section 7
        "improvement_rows": improvement_rows,
        # Summary
        "framework_rows": framework_rows,
    }


# ═══════════════════════════════════════════════════════════════════
# PowerPoint presentation
# ═══════════════════════════════════════════════════════════════════


def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _pptx_set_cell_border(cell, color_tuple):
    """Set thin borders on all four sides of a table cell via XML."""
    from lxml import etree
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = etree.SubElement(tc_pr, f"{{{nsmap['a']}}}{edge}", attrib={"w": "6350", "cmpd": "sng"})
        solid = etree.SubElement(ln, f"{{{nsmap['a']}}}solidFill")
        etree.SubElement(solid, f"{{{nsmap['a']}}}srgbClr", attrib={
            "val": f"{color_tuple[0]:02X}{color_tuple[1]:02X}{color_tuple[2]:02X}",
        })


def _pptx_style_table(table, headers, rows, col_widths=None):
    """Apply consistent styling to a python-pptx table."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Remove the built-in banding
    tbl = table._tbl
    tbl_pr = tbl.tblPr
    tbl_pr.set("bandRow", "0")
    tbl_pr.set("firstRow", "0")
    tbl_pr.set("lastRow", "0")

    n_rows = len(rows)

    # Header row
    for c_idx, hdr in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(*_CLR_HEADER_BG)
        _pptx_set_cell_border(cell, _CLR_BORDER)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = "Calibri"
            p.font.color.rgb = _rgb(*_CLR_HEADER_FG)
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(1)
            p.space_after = Pt(1)

    # Data rows
    for r_idx in range(n_rows):
        for c_idx in range(len(headers)):
            cell = table.cell(r_idx + 1, c_idx)
            val = rows[r_idx][c_idx] if c_idx < len(rows[r_idx]) else ""
            cell.text = str(val) if val else ""
            _pptx_set_cell_border(cell, _CLR_BORDER)
            # Alternating fill
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(*_CLR_ALT_ROW)
            else:
                cell.fill.background()
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.name = "Calibri"
                p.font.color.rgb = _rgb(*_CLR_DARK)
                p.space_before = Pt(1)
                p.space_after = Pt(1)


def _pptx_add_table_slide(prs, title, headers, rows, col_widths=None):
    """Add a slide with a title and a styled data table."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title text box
    txb = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(12), Inches(0.55),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = _rgb(*_CLR_HEADING)
    p.alignment = PP_ALIGN.LEFT

    # Thin accent line below the title
    slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.6), Inches(0.88), Inches(12), Inches(0.02),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _rgb(*_CLR_ACCENT)
    slide.shapes[-1].line.fill.background()

    # Table
    capped_rows = rows[:15]
    n_data = len(capped_rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_data + 1, n_cols,
        Inches(0.6), Inches(1.05),
        Inches(12), Inches(0.28 * (n_data + 1)),
    )
    _pptx_style_table(table_shape.table, headers, capped_rows, col_widths)

    return slide


def _pptx_add_content_slide(prs, title, bullet_points):
    """Add a slide with title, accent line, and bullet list."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    txb = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.3), Inches(12), Inches(0.55),
    )
    p = txb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = _rgb(*_CLR_HEADING)

    # Accent line
    slide.shapes.add_shape(
        1, Inches(0.6), Inches(0.88), Inches(12), Inches(0.02),
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _rgb(*_CLR_ACCENT)
    slide.shapes[-1].line.fill.background()

    # Bullet body
    txb = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.15), Inches(11.6), Inches(5.8),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Indent sub-items that start with spaces
        is_sub = text.startswith("  ")
        p.text = text.strip()
        p.font.size = Pt(14) if not is_sub else Pt(12)
        p.font.name = "Calibri"
        p.font.color.rgb = _rgb(*_CLR_DARK) if not is_sub else _rgb(*_CLR_MUTED)
        p.space_before = Pt(6)
        p.space_after = Pt(2)
        if is_sub:
            p.level = 1

    return slide


def _pptx_add_section_slide(prs, title):
    """Add a visually distinct section divider slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Full-width accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(3.0), Inches(13.333), Inches(0.04),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_CLR_ACCENT)
    bar.line.fill.background()

    # Section title centred
    txb = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(10.333), Inches(2.5),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = _rgb(*_CLR_HEADING)
    p.alignment = PP_ALIGN.CENTER

    return slide


def _pptx_add_title_slide(prs, title, subtitle):
    """Add a clean title slide with centred text."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Top accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_CLR_ACCENT)
    bar.line.fill.background()

    # Main title
    txb = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.5),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = _rgb(*_CLR_HEADING)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txb = slide.shapes.add_textbox(
        Inches(2.0), Inches(3.6), Inches(9.333), Inches(1.0),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.name = "Calibri"
    p.font.color.rgb = _rgb(*_CLR_MUTED)
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_CLR_ACCENT)
    bar.line.fill.background()

    return slide


def _format_period_label(period_start, period_end):
    """Return a human-readable period label for the cover page."""
    if period_start and period_end:
        return f"Periode : {period_start.strftime('%d/%m/%Y')} - {period_end.strftime('%d/%m/%Y')}"
    if period_start:
        return f"Depuis le {period_start.strftime('%d/%m/%Y')}"
    if period_end:
        return f"Jusqu'au {period_end.strftime('%d/%m/%Y')}"
    return ""


def generate_management_review_pptx(user, scope_ids=None,
                                     period_start=None, period_end=None):
    """Generate a management review PowerPoint presentation.

    Returns a tuple (filename, content_bytes).
    """
    from pptx import Presentation
    from pptx.util import Inches

    data = gather_management_review_data(
        user, scope_ids,
        period_start=period_start, period_end=period_end,
    )
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    company_name = data["company"].name or "Fairway"
    date_str = data["generated_at"].strftime("%d/%m/%Y")
    period_label = _format_period_label(data["period_start"], data["period_end"])

    subtitle = f"{company_name}  -  {date_str}"
    if period_label:
        subtitle += f"\n{period_label}"

    # ── Title slide ──
    _pptx_add_title_slide(
        prs,
        "Revue de direction - SMSI",
        subtitle,
    )

    # ── Agenda ──
    _pptx_add_content_slide(prs, "Ordre du jour", [
        "1.  Etat d'avancement des actions decidees lors de revues precedentes",
        "2.  Modifications des enjeux externes et internes",
        "3.  Modifications des besoins et attentes des parties interessees",
        "4.  Retours sur les performances SSI",
        "5.  Retour d'informations des parties interessees",
        "6.  Resultats de l'appreciation des risques et plan de traitement",
        "7.  Opportunites d'amelioration",
    ])

    # ── 1. Action plan status ──
    _pptx_add_section_slide(prs, "1.  Etat d'avancement des actions")

    stats = data["action_plan_stats"]
    bullets = [f"Total des plans d'actions : {stats['total']}"]
    for label, count in stats["by_status"].items():
        bullets.append(f"  {label} : {count}")
    if stats["overdue"]:
        bullets.append(f"En retard : {stats['overdue']}")
    _pptx_add_content_slide(prs, "Synthese des plans d'actions", bullets)

    if data["action_plan_rows"]:
        _pptx_add_table_slide(
            prs, "Plans d'actions en cours",
            ["Ref.", "Nom", "Statut", "Responsable", "Echeance", "Avancement"],
            [[r["reference"], r["name"][:50], r["status"], r["owner"],
              r["target_date"], r["progress"]]
             for r in data["action_plan_rows"]],
            col_widths=[1.0, 4.0, 1.6, 2.0, 1.4, 1.2],
        )

    # ── 2. Issues ──
    _pptx_add_section_slide(prs, "2.  Enjeux externes et internes")

    if data["internal_issues"]:
        _pptx_add_table_slide(
            prs, "Enjeux internes",
            ["Ref.", "Enjeu", "Categorie", "Impact", "Tendance", "Statut"],
            [[i["reference"], i["name"][:50], i["category"],
              i["impact_level"], i["trend"], i["status"]]
             for i in data["internal_issues"]],
            col_widths=[1.0, 4.0, 1.8, 1.4, 1.4, 1.4],
        )

    if data["external_issues"]:
        _pptx_add_table_slide(
            prs, "Enjeux externes",
            ["Ref.", "Enjeu", "Categorie", "Impact", "Tendance", "Statut"],
            [[i["reference"], i["name"][:50], i["category"],
              i["impact_level"], i["trend"], i["status"]]
             for i in data["external_issues"]],
            col_widths=[1.0, 4.0, 1.8, 1.4, 1.4, 1.4],
        )

    # ── 3. Stakeholders ──
    _pptx_add_section_slide(prs, "3.  Besoins et attentes des parties interessees")

    if data["stakeholder_rows"]:
        _pptx_add_table_slide(
            prs, "Parties interessees",
            ["Ref.", "Nom", "Type", "Categorie", "Influence", "Interet", "Attentes"],
            [[s["reference"], s["name"][:35], s["type"], s["category"],
              s["influence_level"], s["interest_level"],
              str(s["expectations_count"])]
             for s in data["stakeholder_rows"]],
            col_widths=[1.0, 3.0, 1.2, 2.0, 1.4, 1.4, 1.0],
        )

    # ── 4. Information security performance ──
    _pptx_add_section_slide(prs, "4.  Retours sur les performances SSI")

    # 4a - NC
    if data["finding_type_counts"]:
        bullets = ["Repartition des constats :"]
        for label, count in data["finding_type_counts"].items():
            bullets.append(f"  {label} : {count}")
        _pptx_add_content_slide(prs, "4a.  Non-conformites et actions correctives", bullets)

    if data["nc_rows"]:
        _pptx_add_table_slide(
            prs, "Non-conformites",
            ["Ref.", "Type", "Description", "Audit", "Recommandation"],
            [[n["reference"], n["type"], n["description"][:80],
              n["assessment"][:35], n["recommendation"][:80]]
             for n in data["nc_rows"]],
            col_widths=[1.0, 1.8, 3.5, 2.2, 3.5],
        )

    # 4b - Indicators
    if data["indicator_rows"]:
        _pptx_add_table_slide(
            prs, "4b.  Resultats de surveillance et mesurage",
            ["Ref.", "Indicateur", "Valeur", "Cible", "Unite", "Critique", "Frequence"],
            [[i["reference"], i["name"][:35], i["current_value"],
              i["expected_level"], i["unit"],
              "OUI" if i["is_critical"] else "Non",
              i["review_frequency"]]
             for i in data["indicator_rows"]],
            col_widths=[1.0, 3.2, 1.4, 1.4, 1.0, 1.0, 1.4],
        )

    # 4c - Audit results
    if data["assessment_rows"]:
        _pptx_add_table_slide(
            prs, "4c.  Resultats des audits",
            ["Ref.", "Nom", "Referentiels", "Conformite", "Couverture", "Date fin"],
            [[a["reference"], a["name"][:35], a["frameworks"][:35],
              a["compliance_level"], a["coverage"], a["end_date"]]
             for a in data["assessment_rows"]],
            col_widths=[1.0, 3.0, 3.0, 1.4, 1.4, 1.4],
        )

    if data["framework_rows"]:
        _pptx_add_table_slide(
            prs, "Synthese de conformite par referentiel",
            ["Referentiel", "Conformite", "Nb exigences", "Statut"],
            [[f["name"], f["compliance_level"],
              str(f["requirement_count"]), f["status"]]
             for f in data["framework_rows"]],
            col_widths=[5.0, 2.0, 2.0, 2.0],
        )

    # 4d - Objectives
    obj_stats = data["objective_stats"]
    if obj_stats["total"]:
        _pptx_add_content_slide(prs, "4d.  Realisation des objectifs SSI", [
            f"Total des objectifs : {obj_stats['total']}",
            f"  Atteints : {obj_stats['achieved']}",
            f"  Non atteints : {obj_stats['not_achieved']}",
            f"  En cours : {obj_stats['in_progress']}",
        ])

    if data["objective_rows"]:
        _pptx_add_table_slide(
            prs, "Objectifs de securite",
            ["Ref.", "Objectif", "Categorie", "Statut", "Avancement", "Responsable", "Echeance"],
            [[o["reference"], o["name"][:35], o["category"],
              o["status"], o["progress"], o["owner"], o["target_date"]]
             for o in data["objective_rows"]],
            col_widths=[1.0, 3.0, 1.6, 1.4, 1.2, 2.0, 1.2],
        )

    # ── 5. Feedback from interested parties ──
    _pptx_add_section_slide(prs, "5.  Retour d'informations des parties interessees")

    if data["expectation_rows"]:
        _pptx_add_table_slide(
            prs, "Attentes des parties interessees",
            ["Partie interesssee", "Description", "Type", "Priorite"],
            [[e["stakeholder"][:30], e["description"][:80],
              e["type"], e["priority"]]
             for e in data["expectation_rows"]],
            col_widths=[2.5, 6.0, 1.5, 1.5],
        )

    # ── 6. Risk assessment and treatment ──
    _pptx_add_section_slide(prs, "6.  Risques et plan de traitement")

    r_stats = data["risk_stats"]
    bullets = [f"Total des risques : {r_stats['total']}"]
    for label, count in r_stats["by_treatment"].items():
        bullets.append(f"  {label} : {count}")
    _pptx_add_content_slide(prs, "Synthese des risques", bullets)

    if data["risk_assessment_rows"]:
        _pptx_add_table_slide(
            prs, "Appreciations des risques",
            ["Ref.", "Nom", "Methodologie", "Statut", "Nb risques", "Date"],
            [[ra["reference"], ra["name"][:35], ra["methodology"],
              ra["status"], str(ra["risk_count"]), ra["date"]]
             for ra in data["risk_assessment_rows"]],
            col_widths=[1.0, 3.5, 2.0, 1.6, 1.4, 1.4],
        )

    if data["critical_risk_rows"]:
        _pptx_add_table_slide(
            prs, "Risques critiques / eleves",
            ["Ref.", "Risque", "Niveau", "Traitement", "Statut"],
            [[r["reference"], r["name"][:50], str(r["risk_level"]),
              r["treatment"], r["status"]]
             for r in data["critical_risk_rows"]],
            col_widths=[1.0, 5.0, 1.2, 2.0, 1.6],
        )

    tp = data["tp_stats"]
    bullets = [f"Total des plans de traitement : {tp['total']}"]
    for label, count in tp["by_status"].items():
        bullets.append(f"  {label} : {count}")
    if tp["overdue"]:
        bullets.append(f"En retard : {tp['overdue']}")
    _pptx_add_content_slide(prs, "Plans de traitement des risques", bullets)

    if data["treatment_plan_rows"]:
        _pptx_add_table_slide(
            prs, "Plans de traitement en cours",
            ["Ref.", "Nom", "Risque", "Statut", "Responsable", "Echeance", "Avancement"],
            [[t["reference"], t["name"][:35], t["risk"][:30],
              t["status"], t["owner"], t["target_date"], t["progress"]]
             for t in data["treatment_plan_rows"]],
            col_widths=[1.0, 2.8, 2.5, 1.4, 1.8, 1.2, 1.2],
        )

    # ── 7. Opportunities for improvement ──
    _pptx_add_section_slide(prs, "7.  Opportunites d'amelioration")

    if data["improvement_rows"]:
        _pptx_add_table_slide(
            prs, "Opportunites d'amelioration identifiees",
            ["Ref.", "Description", "Recommandation", "Audit"],
            [[i["reference"], i["description"][:80],
              i["recommendation"][:80], i["assessment"][:35]]
             for i in data["improvement_rows"]],
            col_widths=[1.0, 4.5, 4.5, 2.0],
        )
    else:
        _pptx_add_content_slide(
            prs,
            "Opportunites d'amelioration",
            ["Aucune opportunite d'amelioration identifiee."],
        )

    # Save
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    date_str = data["generated_at"].strftime("%Y%m%d_%H%M%S")
    filename = f"Management_Review_{date_str}.pptx"

    return filename, buf.getvalue()


# ═══════════════════════════════════════════════════════════════════
# Word meeting minutes (DOCX)
# ═══════════════════════════════════════════════════════════════════


def _docx_setup_styles(doc):
    """Configure base document styles for a neutral, readable look."""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_LINE_SPACING

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10)
    style.font.color.rgb = RGBColor(*_CLR_DARK)
    pf = style.paragraph_format
    pf.space_before = Pt(2)
    pf.space_after = Pt(4)
    pf.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    pf.line_spacing = 1.15

    for level in range(1, 4):
        hstyle = doc.styles[f"Heading {level}"]
        hstyle.font.name = "Calibri"
        hstyle.font.color.rgb = RGBColor(*_CLR_HEADING)
        if level == 1:
            hstyle.font.size = Pt(16)
        elif level == 2:
            hstyle.font.size = Pt(13)
        else:
            hstyle.font.size = Pt(11)
        hstyle.font.bold = True
        hpf = hstyle.paragraph_format
        hpf.space_before = Pt(14 if level == 1 else 10)
        hpf.space_after = Pt(6)


def _docx_add_table(doc, headers, rows):
    """Add a clean formatted table with header and data rows."""
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.enum.table import WD_TABLE_ALIGNMENT

    if not rows:
        doc.add_paragraph("Aucun element.", style="List Bullet")
        return None

    table = doc.add_table(rows=len(rows) + 1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    table.autofit = True

    # Apply a clean style then override colours
    table.style = "Table Grid"

    # Header row
    for c_idx, hdr in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = ""
        p = cell.paragraphs[0]
        run = p.add_run(hdr)
        run.bold = True
        run.font.size = Pt(9)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(*_CLR_HEADER_FG)
        # Background
        shading = cell._tc.get_or_add_tcPr()
        shading_el = shading.makeelement(qn("w:shd"), {
            qn("w:fill"): f"{_CLR_HEADER_BG[0]:02X}{_CLR_HEADER_BG[1]:02X}{_CLR_HEADER_BG[2]:02X}",
            qn("w:val"): "clear",
        })
        shading.append(shading_el)

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(str(val) if val else "")
            run.font.size = Pt(9)
            run.font.name = "Calibri"
            run.font.color.rgb = RGBColor(*_CLR_DARK)
            # Alternating row
            if r_idx % 2 == 1:
                shading = cell._tc.get_or_add_tcPr()
                shading_el = shading.makeelement(qn("w:shd"), {
                    qn("w:fill"): f"{_CLR_ALT_ROW[0]:02X}{_CLR_ALT_ROW[1]:02X}{_CLR_ALT_ROW[2]:02X}",
                    qn("w:val"): "clear",
                })
                shading.append(shading_el)

    doc.add_paragraph()  # spacing after table
    return table


def _docx_add_kv_table(doc, pairs):
    """Add a key-value metadata table (2 columns, label + value)."""
    from docx.shared import Pt
    from docx.oxml.ns import qn

    table = doc.add_table(rows=len(pairs), cols=2)
    table.style = "Table Grid"
    table.autofit = True

    for i, (label, value) in enumerate(pairs):
        # Label cell
        lc = table.cell(i, 0)
        lc.text = ""
        run = lc.paragraphs[0].add_run(label)
        run.bold = True
        run.font.size = Pt(10)
        run.font.name = "Calibri"
        shading = lc._tc.get_or_add_tcPr()
        shading_el = shading.makeelement(qn("w:shd"), {
            qn("w:fill"): f"{_CLR_HEADER_BG[0]:02X}{_CLR_HEADER_BG[1]:02X}{_CLR_HEADER_BG[2]:02X}",
            qn("w:val"): "clear",
        })
        shading.append(shading_el)
        # Value cell
        vc = table.cell(i, 1)
        vc.text = ""
        run = vc.paragraphs[0].add_run(value)
        run.font.size = Pt(10)
        run.font.name = "Calibri"

    doc.add_paragraph()
    return table


def _docx_add_decisions_block(doc):
    """Add a standard decisions placeholder after each section."""
    from docx.shared import Pt, RGBColor
    p = doc.add_paragraph()
    run = p.add_run("Decisions :")
    run.bold = True
    run.font.size = Pt(10)
    run.font.name = "Calibri"
    p = doc.add_paragraph("[A completer]", style="List Bullet")
    for run in p.runs:
        run.font.color.rgb = RGBColor(*_CLR_MUTED)


def generate_management_review_docx(user, scope_ids=None,
                                     period_start=None, period_end=None):
    """Generate a management review meeting minutes document (DOCX).

    Returns a tuple (filename, content_bytes).
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    data = gather_management_review_data(
        user, scope_ids,
        period_start=period_start, period_end=period_end,
    )
    doc = Document()
    _docx_setup_styles(doc)

    company_name = data["company"].name or "Fairway"
    date_str = data["generated_at"].strftime("%d/%m/%Y")
    period_label = _format_period_label(data["period_start"], data["period_end"])

    # ── Cover page ──
    for _ in range(4):
        doc.add_paragraph()  # vertical spacing

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(company_name)
    run.font.size = Pt(28)
    run.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(*_CLR_HEADING)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(12)

    # Horizontal rule via bottom border on paragraph
    p_fmt = p._p.get_or_add_pPr()
    p_bdr = p_fmt.makeelement(qn("w:pBdr"), {})
    bottom = p_bdr.makeelement(qn("w:bottom"), {
        qn("w:val"): "single",
        qn("w:sz"): "6",
        qn("w:space"): "1",
        qn("w:color"): f"{_CLR_ACCENT[0]:02X}{_CLR_ACCENT[1]:02X}{_CLR_ACCENT[2]:02X}",
    })
    p_bdr.append(bottom)
    p_fmt.append(p_bdr)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(8)
    run = p.add_run("Compte rendu de revue de direction")
    run.font.size = Pt(20)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(*_CLR_HEADING)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Systeme de management de la securite de l'information")
    run.font.size = Pt(13)
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(*_CLR_MUTED)

    if period_label:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(6)
        run = p.add_run(period_label)
        run.font.size = Pt(12)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(*_CLR_MUTED)

    doc.add_paragraph()
    doc.add_paragraph()

    # Metadata table
    meta_rows = [
        ("Date", date_str),
    ]
    if period_label:
        meta_rows.append(("Periode", period_label))
    meta_rows += [
        ("Lieu", "[A completer]"),
        ("Participants", "[A completer]"),
        ("Redacteur", str(user)),
    ]
    _docx_add_kv_table(doc, meta_rows)

    doc.add_page_break()

    # ── Table of contents ──
    doc.add_heading("Ordre du jour", level=1)
    toc_items = [
        "Etat d'avancement des actions decidees lors de revues precedentes",
        "Modifications des enjeux externes et internes pertinents pour le SMSI",
        "Modifications des besoins et attentes des parties interessees pertinentes pour le SMSI",
        "Retours sur les performances SSI",
        "Retour d'informations des parties interessees",
        "Resultats de l'appreciation des risques et etat d'avancement du plan de traitement des risques",
        "Opportunites d'amelioration",
    ]
    for idx, item in enumerate(toc_items, 1):
        p = doc.add_paragraph()
        run = p.add_run(f"{idx}.  {item}")
        run.font.size = Pt(10)
        run.font.name = "Calibri"
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after = Pt(2)

    doc.add_page_break()

    # ── 1. Action plans ──
    doc.add_heading(
        "1.  Etat d'avancement des actions decidees lors de revues precedentes",
        level=1,
    )

    stats = data["action_plan_stats"]
    doc.add_paragraph(f"Total des plans d'actions : {stats['total']}")
    for label, count in stats["by_status"].items():
        doc.add_paragraph(f"{label} : {count}", style="List Bullet")
    if stats["overdue"]:
        p = doc.add_paragraph()
        run = p.add_run(f"Plans d'actions en retard : {stats['overdue']}")
        run.bold = True

    if data["action_plan_rows"]:
        doc.add_heading("Detail des plans d'actions en cours", level=2)
        _docx_add_table(doc,
            ["Ref.", "Nom", "Statut", "Responsable", "Echeance", "Avancement"],
            [[r["reference"], r["name"], r["status"], r["owner"],
              r["target_date"], r["progress"]]
             for r in data["action_plan_rows"]],
        )

    _docx_add_decisions_block(doc)

    # ── 2. Issues ──
    doc.add_heading(
        "2.  Modifications des enjeux externes et internes pertinents pour le SMSI",
        level=1,
    )

    if data["internal_issues"]:
        doc.add_heading("Enjeux internes", level=2)
        _docx_add_table(doc,
            ["Ref.", "Enjeu", "Categorie", "Impact", "Tendance", "Statut"],
            [[i["reference"], i["name"], i["category"],
              i["impact_level"], i["trend"], i["status"]]
             for i in data["internal_issues"]],
        )

    if data["external_issues"]:
        doc.add_heading("Enjeux externes", level=2)
        _docx_add_table(doc,
            ["Ref.", "Enjeu", "Categorie", "Impact", "Tendance", "Statut"],
            [[i["reference"], i["name"], i["category"],
              i["impact_level"], i["trend"], i["status"]]
             for i in data["external_issues"]],
        )

    _docx_add_decisions_block(doc)

    # ── 3. Stakeholders ──
    doc.add_heading(
        "3.  Modifications des besoins et attentes des parties interessees",
        level=1,
    )

    if data["stakeholder_rows"]:
        _docx_add_table(doc,
            ["Ref.", "Nom", "Type", "Categorie", "Influence", "Interet", "Nb attentes"],
            [[s["reference"], s["name"], s["type"], s["category"],
              s["influence_level"], s["interest_level"],
              str(s["expectations_count"])]
             for s in data["stakeholder_rows"]],
        )

    _docx_add_decisions_block(doc)

    # ── 4. SSI performance ──
    doc.add_heading("4.  Retours sur les performances SSI", level=1)

    # 4a
    doc.add_heading("4a.  Non-conformites et actions correctives", level=2)

    if data["finding_type_counts"]:
        doc.add_paragraph("Repartition des constats :")
        for label, count in data["finding_type_counts"].items():
            doc.add_paragraph(f"{label} : {count}", style="List Bullet")

    if data["nc_rows"]:
        doc.add_heading("Detail des non-conformites", level=3)
        _docx_add_table(doc,
            ["Ref.", "Type", "Description", "Audit", "Recommandation"],
            [[n["reference"], n["type"], n["description"],
              n["assessment"], n["recommendation"]]
             for n in data["nc_rows"]],
        )

    # 4b
    doc.add_heading("4b.  Resultats de la surveillance et du mesurage", level=2)

    if data["indicator_rows"]:
        _docx_add_table(doc,
            ["Ref.", "Indicateur", "Valeur", "Cible", "Unite", "Critique", "Frequence"],
            [[i["reference"], i["name"], i["current_value"],
              i["expected_level"], i["unit"],
              "OUI" if i["is_critical"] else "Non",
              i["review_frequency"]]
             for i in data["indicator_rows"]],
        )
    else:
        doc.add_paragraph("Aucun indicateur actif.", style="List Bullet")

    # 4c
    doc.add_heading("4c.  Resultats des audits", level=2)

    if data["assessment_rows"]:
        _docx_add_table(doc,
            ["Ref.", "Nom", "Referentiels", "Conformite", "Couverture", "Date fin"],
            [[a["reference"], a["name"], a["frameworks"],
              a["compliance_level"], a["coverage"], a["end_date"]]
             for a in data["assessment_rows"]],
        )
    else:
        doc.add_paragraph("Aucun audit termine.", style="List Bullet")

    if data["framework_rows"]:
        doc.add_heading("Conformite par referentiel", level=3)
        _docx_add_table(doc,
            ["Referentiel", "Conformite", "Nb exigences", "Statut"],
            [[f["name"], f["compliance_level"],
              str(f["requirement_count"]), f["status"]]
             for f in data["framework_rows"]],
        )

    # 4d
    doc.add_heading("4d.  Realisation des objectifs SSI", level=2)

    obj_stats = data["objective_stats"]
    if obj_stats["total"]:
        doc.add_paragraph(f"Total des objectifs : {obj_stats['total']}")
        doc.add_paragraph(f"Atteints : {obj_stats['achieved']}", style="List Bullet")
        doc.add_paragraph(f"Non atteints : {obj_stats['not_achieved']}", style="List Bullet")
        doc.add_paragraph(f"En cours : {obj_stats['in_progress']}", style="List Bullet")

    if data["objective_rows"]:
        _docx_add_table(doc,
            ["Ref.", "Objectif", "Categorie", "Statut", "Avancement", "Responsable", "Echeance"],
            [[o["reference"], o["name"], o["category"],
              o["status"], o["progress"], o["owner"], o["target_date"]]
             for o in data["objective_rows"]],
        )

    _docx_add_decisions_block(doc)

    # ── 5. Feedback from interested parties ──
    doc.add_heading(
        "5.  Retour d'informations des parties interessees",
        level=1,
    )

    if data["expectation_rows"]:
        _docx_add_table(doc,
            ["Partie interessee", "Description", "Type", "Priorite"],
            [[e["stakeholder"], e["description"],
              e["type"], e["priority"]]
             for e in data["expectation_rows"]],
        )
    else:
        doc.add_paragraph("Aucune attente applicable identifiee.", style="List Bullet")

    _docx_add_decisions_block(doc)

    # ── 6. Risks and treatment ──
    doc.add_heading(
        "6.  Resultats de l'appreciation des risques et plan de traitement",
        level=1,
    )

    r_stats = data["risk_stats"]
    doc.add_paragraph(f"Total des risques : {r_stats['total']}")
    for label, count in r_stats["by_treatment"].items():
        doc.add_paragraph(f"{label} : {count}", style="List Bullet")

    if data["risk_assessment_rows"]:
        doc.add_heading("Appreciations des risques", level=2)
        _docx_add_table(doc,
            ["Ref.", "Nom", "Methodologie", "Statut", "Nb risques", "Date"],
            [[ra["reference"], ra["name"], ra["methodology"],
              ra["status"], str(ra["risk_count"]), ra["date"]]
             for ra in data["risk_assessment_rows"]],
        )

    if data["critical_risk_rows"]:
        doc.add_heading("Risques critiques / eleves", level=2)
        _docx_add_table(doc,
            ["Ref.", "Risque", "Niveau", "Traitement", "Statut"],
            [[r["reference"], r["name"], str(r["risk_level"]),
              r["treatment"], r["status"]]
             for r in data["critical_risk_rows"]],
        )

    doc.add_heading("Plans de traitement des risques", level=2)
    tp = data["tp_stats"]
    doc.add_paragraph(f"Total : {tp['total']}")
    for label, count in tp["by_status"].items():
        doc.add_paragraph(f"{label} : {count}", style="List Bullet")
    if tp["overdue"]:
        p = doc.add_paragraph()
        run = p.add_run(f"En retard : {tp['overdue']}")
        run.bold = True

    if data["treatment_plan_rows"]:
        _docx_add_table(doc,
            ["Ref.", "Nom", "Risque", "Statut", "Responsable", "Echeance", "Avancement"],
            [[t["reference"], t["name"], t["risk"],
              t["status"], t["owner"], t["target_date"], t["progress"]]
             for t in data["treatment_plan_rows"]],
        )

    _docx_add_decisions_block(doc)

    # ── 7. Improvement opportunities ──
    doc.add_heading("7.  Opportunites d'amelioration", level=1)

    if data["improvement_rows"]:
        _docx_add_table(doc,
            ["Ref.", "Description", "Recommandation", "Audit"],
            [[i["reference"], i["description"],
              i["recommendation"], i["assessment"]]
             for i in data["improvement_rows"]],
        )
    else:
        doc.add_paragraph(
            "Aucune opportunite d'amelioration identifiee.", style="List Bullet",
        )

    _docx_add_decisions_block(doc)

    # ── Closing section ──
    doc.add_page_break()
    doc.add_heading("Synthese des decisions", level=1)
    p = doc.add_paragraph("[A completer]")
    for run in p.runs:
        run.font.color.rgb = RGBColor(*_CLR_MUTED)

    doc.add_heading("Prochaine revue de direction", level=2)
    p = doc.add_paragraph("Date prevue : [A completer]")
    for run in p.runs:
        run.font.color.rgb = RGBColor(*_CLR_MUTED)

    doc.add_heading("Signatures", level=2)
    _docx_add_table(doc,
        ["Nom", "Fonction", "Signature"],
        [["", "", ""], ["", "", ""]],
    )

    # Save
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    date_str = data["generated_at"].strftime("%Y%m%d_%H%M%S")
    filename = f"Management_Review_Minutes_{date_str}.docx"

    return filename, buf.getvalue()
